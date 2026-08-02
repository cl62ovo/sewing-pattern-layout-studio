"""Generate an editable Plush Pattern Studio pitch deck with python-pptx.

Run: python presentation/build_pptx.py
Output: presentation/slides.pptx (native, editable text boxes — not an image export)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ACCENT = RGBColor(0xB5, 0x48, 0x2A)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x77, 0x77, 0x77)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Demo order swapped per request: Nest & Cut first, then Pattern Studio walkthrough.
SLIDES = [
    {
        "type": "title",
        "title": "Plush Pattern Studio",
        "subtitle": "From an idea to a sewable pattern — powered by AI\n\nTeam: [Name A] & [Name B]",
    },
    {
        "type": "paragraphs",
        "title": "The Problem",
        "paragraphs": [
            [("Imagine you want to make a plush doll from scratch.", False)],
            [(
                "First, you need a technical 2D blueprint, complete with darts and seam "
                "allowances. You have to cut out each paper pattern piece, lay them "
                "strategically on a sheet of fabric to maximize material, trace every "
                "outline onto the cloth, and then carefully cut out all the fabric pieces. "
                "Only then can you start sewing them together into a 3D toy.",
                False,
            )],
            [
                ("As you can see, this preparation requires both ", False),
                ("geometric thinking", True),
                (" and ", False),
                ("tedious manual work", True),
                (" — and it hits two very different groups hard.", False),
            ],
        ],
    },
    {
        "type": "sections",
        "title": "Who Feels the Pain",
        "sections": [
            {
                "heading": "For Professionals — the cost of complexity",
                "bullets": [
                    "Inefficient layouts: manual nesting is a slow, error-prone \"geometric puzzle\"",
                    "High fabric waste: poor nesting wastes 15%\u201330% of expensive material",
                    "Alignment failure: matching stripes/checks by hand is unforgiving — one mistake ruins the garment",
                ],
            },
            {
                "heading": "For Beginners — the barrier to entry",
                "bullets": [
                    "The \"Skill Wall\": no pattern-making knowledge (darts, seam allowances, notches) keeps ideas stuck on paper",
                    "Complexity paralysis: math, tool choices, and dense tutorials overwhelm novices before they start",
                    "Low success rate: frustration kills the joy of handmade creation",
                ],
            },
        ],
    },
    {
        "type": "solution",
        "title": "Our Solution",
        "intro": [
            ("That's why we created ", False),
            ("Plush Pattern Studio", True),
            (" — to take care of this complex preparation so anyone can jump straight to the joy of creating.", False),
        ],
        "bullets": [
            "Describe your plush doll in plain language + pick a target height",
            "AI turns your words into a structured design spec",
            "A 3D model candidate is generated for you to confirm or regenerate",
            "The system auto-normalizes the mesh, plans seams, and unfolds panels",
            "Output: orthographic views + a 1:1 A4 sewing pattern, ready to print and cut",
        ],
    },
    {
        "type": "placeholder",
        "title": "Live Demo — For Professionals: Nest & Cut",
        "note": "(demo placeholder — to be filled in)",
    },
    {
        "type": "bullets",
        "title": "Live Demo — From Words to Pattern",
        "bullets": [
            "Type a description + choose a height",
            "AI drafts the design → a 3D model is generated",
            "Preview and accept it (or regenerate with edits)",
            "Get 2D pattern pieces + a printable A4 PDF",
        ],
        "note": "[占位符：demo 截图 / GIF]",
    },
    {
        "type": "sections",
        "title": "How We Built It",
        "sections": [
            {
                "heading": "v1 — Vibe coding, fast iteration",
                "bullets": [
                    "Talked through requirements in a Q&A dialogue → wrote them into a design doc",
                    "Vibe-coded a React MVP with Codex, then iterated",
                ],
            },
            {
                "heading": "v2 — Adding the real backend",
                "bullets": [
                    "FastAPI backend + PostgreSQL + Redis task queue",
                    "OpenRouter for structured requirement parsing",
                    "Meshy API for 3D model generation",
                    "Python geometry worker: mesh repair, normalization, seam planning, unfolding",
                ],
            },
        ],
    },
    {
        "type": "bullets",
        "title": "What's Next",
        "bullets": [
            "Paper size input is still manual — needs smart presets (A4 / Letter / roll widths)",
            "The UI flow is long — needs a smoother, more guided experience",
            "More garment shapes & fabric types",
            "A Nest & Cut layout optimizer for professionals",
        ],
    },
    {
        "type": "title",
        "title": "Thank You",
        "subtitle": "Questions?\n\nPlush Pattern Studio",
    },
]


def add_title(slide, text, top=Inches(0.4), size=32):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    return box


def add_body_box(slide, top=Inches(1.5), height=None):
    height = height or (SLIDE_H - top - Inches(0.5))
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), height)
    box.text_frame.word_wrap = True
    return box


def write_runs(paragraph, runs):
    for i, (text, bold) in enumerate(runs):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(20)
        run.font.bold = bold
        run.font.color.rgb = DARK


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank_layout)

        if spec["type"] == "title":
            box = slide.shapes.add_textbox(Inches(1), Inches(2.5), SLIDE_W - Inches(2), Inches(1.2))
            p = box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = spec["title"]
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = ACCENT

            sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), SLIDE_W - Inches(2), Inches(2))
            sub_box.text_frame.word_wrap = True
            for i, line in enumerate(spec["subtitle"].split("\n")):
                sp = sub_box.text_frame.paragraphs[0] if i == 0 else sub_box.text_frame.add_paragraph()
                sp.alignment = PP_ALIGN.CENTER
                run = sp.add_run()
                run.text = line
                run.font.size = Pt(20)
                run.font.color.rgb = GRAY
            continue

        add_title(slide, spec["title"])
        body = add_body_box(slide)
        tf = body.text_frame
        first = True

        if spec["type"] == "paragraphs":
            for para in spec["paragraphs"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(14)
                write_runs(p, para)

        elif spec["type"] == "bullets":
            for item in spec["bullets"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(8)
                write_runs(p, [("• " + item, False)])
            if spec.get("note"):
                p = tf.add_paragraph()
                p.space_before = Pt(16)
                write_runs(p, [(spec["note"], False)])
                p.runs[0].font.italic = True
                p.runs[0].font.color.rgb = GRAY

        elif spec["type"] == "sections":
            for section in spec["sections"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                write_runs(p, [(section["heading"], True)])
                for item in section["bullets"]:
                    bp = tf.add_paragraph()
                    bp.level = 1
                    bp.space_after = Pt(6)
                    write_runs(bp, [("• " + item, False)])

        elif spec["type"] == "solution":
            p = tf.paragraphs[0]
            first = False
            p.space_after = Pt(18)
            write_runs(p, spec["intro"])
            for item in spec["bullets"]:
                bp = tf.add_paragraph()
                bp.space_after = Pt(8)
                write_runs(bp, [("• " + item, False)])

        elif spec["type"] == "placeholder":
            p = tf.paragraphs[0]
            write_runs(p, [(spec["note"], False)])
            p.runs[0].font.italic = True
            p.runs[0].font.color.rgb = GRAY
            p.runs[0].font.size = Pt(24)

    out_path = __file__.replace("build_pptx.py", "slides.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    build()
