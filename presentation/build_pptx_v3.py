"""Restyle the v2 deck into a polished v3 deck (visual only — no content changes).

All text is read back out of the v2 file and re-emitted verbatim, so wording can
never drift. Everything stays as native, editable PowerPoint shapes.

Run: python presentation/build_pptx_v3.py
In : presentation/Sewing-pattern-layout-studio-slides-v2.pptx
Out: presentation/Sewing-pattern-layout-studio-slides-v3.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
SRC = HERE / "Sewing-pattern-layout-studio-slides-v2.pptx"
OUT = HERE / "Sewing-pattern-layout-studio-slides-v3.pptx"

ACCENT = RGBColor(0xB5, 0x48, 0x2A)
NAVY = RGBColor(0x1F, 0x49, 0x7D)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
HAIRLINE = RGBColor(0xD8, 0xD8, 0xD8)
CARD_BG = RGBColor(0xFA, 0xF6, 0xF3)
CARD_BG_ALT = RGBColor(0xF4, 0xF6, 0xFA)
FONT = "Segoe UI"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- shape helpers -----------------------------------------------------------

def find(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"shape {name!r} not found on slide")


def drop(shape):
    shape._element.getparent().remove(shape._element)


def send_to_back(slide, shape):
    tree = slide.shapes._spTree
    tree.remove(shape._element)
    tree.insert(2, shape._element)


def move_behind(slide, shape, reference):
    slide.shapes._spTree.remove(shape._element)
    reference._element.addprevious(shape._element)


def read_paragraphs(shape):
    """[(level, text)] exactly as authored in the source deck."""
    return [
        (p.level, "".join(r.text for r in p.runs))
        for p in shape.text_frame.paragraphs
    ]


def read_text(shape):
    return "\n".join(t for _, t in read_paragraphs(shape))


def rule(slide, left, top, width, color=ACCENT, thickness=Pt(3)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def card(slide, left, top, width, height, fill=CARD_BG, line=HAIRLINE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


def label_backdrop(slide, caption, size=11.25, pad=0.13):
    """White plate so a caption stays legible on top of artwork."""
    text = read_text(caption)
    plate = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        caption.left - Inches(pad),
        caption.top - Inches(0.03),
        Inches(len(text) * size * 0.0078 + pad * 2),
        caption.height + Inches(0.06),
    )
    plate.adjustments[0] = 0.3
    plate.fill.solid()
    plate.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    plate.line.fill.background()
    plate.shadow.inherit = False
    move_behind(slide, plate, caption)
    return plate


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def style(run, size, bold=False, color=DARK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def para(tf, index, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_after=0, space_before=0, line_spacing=1.15, level=0):
    p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    style(run, size, bold, color)
    return p


def bullet_para(tf, index, text, size, color=DARK, dot=ACCENT, space_after=7,
                line_spacing=1.15):
    """Re-emits ``text`` verbatim, but colours a leading bullet glyph."""
    p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if text.startswith("• "):
        head = p.add_run()
        head.text = "•"
        style(head, size, True, dot)
        body = p.add_run()
        body.text = " " + text[2:]
        style(body, size, False, color)
    else:
        run = p.add_run()
        run.text = text
        style(run, size, False, color)
    return p


def section_title(slide, text, size=32, top=Inches(0.375), left=Inches(0.6),
                  width=Inches(12.13), rule_top=None, rule_width=Inches(1.0)):
    tf = textbox(slide, left, top, width, Inches(0.75))
    para(tf, 0, text, size, bold=True, color=ACCENT, line_spacing=1.0)
    rule(slide, left, rule_top if rule_top is not None else top + Inches(0.72), rule_width)
    return tf


# --- per-slide restyling -----------------------------------------------------

def polish_cover(slide):
    """Slide 1 — align the left column, lift the caption off the artwork."""
    for name in ("TextBox 1", "文本框 11"):
        find(slide, name).left = Inches(0.771)
    label_backdrop(slide, find(slide, "pattern-label"))

    authors = find(slide, "文本框 11")
    for paragraph in authors.text_frame.paragraphs:
        paragraph.line_spacing = 1.25
        for run in paragraph.runs:
            style(run, 16, color=GRAY)


def polish_process(slide):
    """Slide 2 — same caption treatment as the cover."""
    label_backdrop(slide, find(slide, "process-pattern-caption"))


def restyle_demo_nest_and_cut(slide):
    """Slide 5 — title column on the left, video flush to the right edge."""
    title = find(slide, "TextBox 1")
    title_text = read_text(title)
    drop(title)

    video = next(s for s in slide.shapes if s.shape_type == 16)
    video.width = Inches(5.625)
    video.height = SLIDE_H
    video.top = 0
    video.left = SLIDE_W - video.width

    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W - video.width, SLIDE_H
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = CARD_BG
    panel.line.fill.background()
    panel.shadow.inherit = False
    send_to_back(slide, panel)

    rule(slide, Inches(0.85), Inches(2.55), Inches(1.1))
    tf = textbox(slide, Inches(0.85), Inches(2.95), Inches(5.9), Inches(2.4))
    para(tf, 0, title_text, 34, bold=True, color=ACCENT, line_spacing=1.1)


def restyle_demo_words_to_pattern(slide):
    """Slide 6 — compact header + bullets in the clear band above the artwork."""
    title = find(slide, "TextBox 1")
    left_box = find(slide, "TextBox 2")
    right_box = find(slide, "文本框 27")

    title_text = read_text(title)
    items = [t for _, t in read_paragraphs(left_box) + read_paragraphs(right_box) if t]
    for shape in (title, left_box, right_box):
        drop(shape)

    # The stacked screenshots differ by a few thousandths of an inch, which reads
    # as jitter when clicking through them — snap them onto one frame.
    for shape in slide.shapes:
        if shape.width and shape.width > Inches(9):
            shape.left = Inches(1.68)
            shape.top = Inches(2.42)
            shape.width = Inches(9.97)
            shape.height = Inches(4.89)

    tf = textbox(slide, Inches(0.6), Inches(0.3), Inches(6.5), Inches(0.6))
    para(tf, 0, title_text, 22, bold=True, color=ACCENT, line_spacing=1.0)
    rule(slide, Inches(0.6), Inches(0.86), Inches(0.9))

    tf = textbox(slide, Inches(0.6), Inches(1.06), Inches(6.0), Inches(1.35))
    for i, item in enumerate(items):
        bullet_para(tf, i, item, 12, space_after=5, line_spacing=1.05)


def restyle_how_we_built_it(slide):
    """Slide 7 — two milestone columns between the screenshots."""
    title = find(slide, "TextBox 1")
    body = find(slide, "TextBox 2")
    title_text = read_text(title)
    paragraphs = [p for p in read_paragraphs(body) if p[1]]
    for shape in (title, body):
        drop(shape)

    groups = []
    for level, text in paragraphs:
        if level == 0:
            groups.append({"heading": text, "bullets": []})
        elif groups:
            groups[-1]["bullets"].append(text)

    section_title(slide, title_text, rule_top=Inches(1.06))

    colors = [NAVY, ACCENT]
    fills = [CARD_BG_ALT, CARD_BG]
    left_positions = [Inches(0.6), Inches(5.5)]
    col_w = Inches(4.6)
    for i, group in enumerate(groups[:2]):
        color = colors[i % 2]
        left = left_positions[i % 2]
        card(slide, left, Inches(1.35), col_w, Inches(3.15), fill=fills[i % 2])
        tf = textbox(slide, left + Inches(0.28), Inches(1.58), col_w - Inches(0.56), Inches(0.4))
        para(tf, 0, group["heading"], 15, bold=True, color=color, line_spacing=1.0)
        rule(slide, left + Inches(0.28), Inches(2.02), Inches(0.7), color, Pt(2.5))
        tf = textbox(slide, left + Inches(0.28), Inches(2.22), col_w - Inches(0.56), Inches(2.1))
        for j, item in enumerate(group["bullets"]):
            bullet_para(tf, j, item, 11.5, dot=color, space_after=6, line_spacing=1.1)


def restyle_whats_next(slide):
    """Slide 8 — one long list split into two balanced cards."""
    title = find(slide, "TextBox 1")
    body = find(slide, "TextBox 2")
    title_text = read_text(title)
    items = [t for _, t in read_paragraphs(body) if t.strip()]
    for shape in (title, body):
        drop(shape)

    section_title(slide, title_text, rule_top=Inches(1.06))

    half = (len(items) + 1) // 2
    columns = [items[:half], items[half:]]
    left_positions = [Inches(0.6), Inches(6.93)]
    col_w = Inches(5.8)
    accents = [ACCENT, NAVY]
    fills = [CARD_BG, CARD_BG_ALT]
    for i, column in enumerate(columns):
        left = left_positions[i]
        card(slide, left, Inches(1.4), col_w, Inches(5.25), fill=fills[i])
        rule(slide, left + Inches(0.45), Inches(1.85), Inches(0.6), accents[i], Pt(2.5))
        tf = textbox(slide, left + Inches(0.45), Inches(2.15), col_w - Inches(0.9),
                     Inches(4.2), anchor=MSO_ANCHOR.MIDDLE)
        for j, item in enumerate(column):
            bullet_para(tf, j, item, 14.5, dot=accents[i], space_after=16, line_spacing=1.2)


def restyle_thank_you(slide):
    """Slide 9 — centred closing card."""
    title = find(slide, "TextBox 1")
    subtitle = find(slide, "TextBox 2")
    title_text = read_text(title)
    subtitle_text = read_text(subtitle)
    for shape in (title, subtitle):
        drop(shape)

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SLIDE_W, Inches(2.4)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = CARD_BG
    band.line.fill.background()
    band.shadow.inherit = False
    send_to_back(slide, band)

    tf = textbox(slide, Inches(1.0), Inches(3.0), SLIDE_W - Inches(2.0), Inches(1.1))
    para(tf, 0, title_text, 52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
         line_spacing=1.0)

    rule(slide, (SLIDE_W - Inches(1.2)) // 2, Inches(4.12), Inches(1.2))

    tf = textbox(slide, Inches(1.0), Inches(4.35), SLIDE_W - Inches(2.0), Inches(0.5))
    para(tf, 0, subtitle_text, 20, color=GRAY, align=PP_ALIGN.CENTER)


# --- deck-wide polish --------------------------------------------------------

def add_top_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def unify_fonts(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT


def build():
    prs = Presentation(str(SRC))

    polish_cover(prs.slides[0])
    polish_process(prs.slides[1])
    restyle_demo_nest_and_cut(prs.slides[4])
    restyle_demo_words_to_pattern(prs.slides[5])
    restyle_how_we_built_it(prs.slides[6])
    restyle_whats_next(prs.slides[7])
    restyle_thank_you(prs.slides[8])

    for slide in prs.slides:
        add_top_bar(slide)
        unify_fonts(slide)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
